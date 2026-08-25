"""在办公文档中执行原文替换，处理跨 run 拆词、页眉页脚、文本框。"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Dict, List


def apply_to_text(text: str, replace_map: Dict[str, str]) -> str:
    """在纯文本上替换；先精确匹配，再允许字之间有空白。"""
    if not text or not replace_map:
        return text
    items = sorted(replace_map.items(), key=lambda kv: len(kv[0]), reverse=True)
    result = text
    for original, token in items:
        if not original:
            continue
        if original in result:
            result = result.replace(original, token)
            continue
        pattern = r"\s*".join(re.escape(ch) for ch in original)
        result = re.sub(pattern, lambda _m, t=token: t, result)
    return result


def _replace_wml_paragraphs(root, replace_map: Dict[str, str]) -> int:
    """在 WordprocessingML 树中按段落拼接 w:t 后替换。返回替换段落数。"""
    from docx.oxml.ns import qn

    changed = 0
    w_p = qn("w:p")
    w_t = qn("w:t")
    xml_space = "{http://www.w3.org/XML/1998/namespace}space"
    for para in root.iter(w_p):
        t_nodes = list(para.iter(w_t))
        if not t_nodes:
            continue
        full = "".join(t.text or "" for t in t_nodes)
        new_full = apply_to_text(full, replace_map)
        if new_full == full:
            continue
        t_nodes[0].text = new_full
        t_nodes[0].set(xml_space, "preserve")
        for t in t_nodes[1:]:
            t.text = ""
        changed += 1
    return changed


def mask_docx(input_path: Path, output_path: Path, replace_map: Dict[str, str]) -> None:
    from docx import Document

    doc = Document(str(input_path))
    _replace_wml_paragraphs(doc.element, replace_map)
    for section in doc.sections:
        for part in (
            section.header,
            section.footer,
            getattr(section, "first_page_header", None),
            getattr(section, "first_page_footer", None),
            getattr(section, "even_page_header", None),
            getattr(section, "even_page_footer", None),
        ):
            if part is None:
                continue
            try:
                _replace_wml_paragraphs(part._element, replace_map)
            except Exception:
                continue
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))


def mask_xlsx(input_path: Path, output_path: Path, replace_map: Dict[str, str]) -> None:
    from openpyxl import load_workbook

    wb = load_workbook(str(input_path))
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is None:
                    continue
                original = str(cell.value)
                new_val = apply_to_text(original, replace_map)
                if new_val != original:
                    cell.value = new_val
    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output_path))


def _iter_pptx_paragraphs(prs) -> List:
    paras = []
    for slide in prs.slides:
        for shape in slide.shapes:
            paras.extend(_paragraphs_from_shape(shape))
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                paras.extend(slide.notes_slide.notes_text_frame.paragraphs)
        except Exception:
            pass
    return paras


def _paragraphs_from_shape(shape) -> List:
    paras = []
    if getattr(shape, "has_text_frame", False):
        paras.extend(shape.text_frame.paragraphs)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    paras.extend(cell.text_frame.paragraphs)
    # 组合形状
    if getattr(shape, "shapes", None):
        try:
            for sub in shape.shapes:
                paras.extend(_paragraphs_from_shape(sub))
        except Exception:
            pass
    return paras


def _replace_pptx_paragraph(paragraph, replace_map: Dict[str, str]) -> None:
    runs = list(paragraph.runs)
    if not runs:
        return
    full = "".join(run.text or "" for run in runs)
    new_full = apply_to_text(full, replace_map)
    if new_full == full:
        return
    runs[0].text = new_full
    for run in runs[1:]:
        run.text = ""


def mask_pptx(input_path: Path, output_path: Path, replace_map: Dict[str, str]) -> None:
    from pptx import Presentation

    prs = Presentation(str(input_path))
    for para in _iter_pptx_paragraphs(prs):
        _replace_pptx_paragraph(para, replace_map)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _pdf_fontname(text: str) -> str:
    if any("\u4e00" <= ch <= "\u9fff" for ch in text):
        return "china-ss"
    return "helv"


def _pdf_fontsize_fit(rect, text: str, fontname: str) -> float:
    import pymupdf as fitz

    height_cap = min(float(rect.height) * 0.9, 12.0)
    lo, hi = 3.0, max(height_cap, 3.0)
    best = lo
    width_limit = max(float(rect.width) - 0.4, 1.0)
    for _ in range(16):
        mid = (lo + hi) / 2
        try:
            width = fitz.get_text_length(text, fontname=fontname, fontsize=mid)
        except Exception:
            return max(4.0, min(float(rect.height) * 0.7, 9.0))
        if width <= width_limit:
            best = mid
            lo = mid
        else:
            hi = mid
    return best


def _pdf_search_rects(page, needle: str):
    """在一页上定位 needle；search_for 失败时按词拼接回退。"""
    import pymupdf as fitz

    if not needle:
        return []
    hits = list(page.search_for(needle) or [])
    if hits:
        return hits
    collapsed = re.sub(r"[\s\u00a0\u200b\u200c\u200d\ufeff]+", "", needle)
    if collapsed and collapsed != needle:
        hits = list(page.search_for(collapsed) or [])
        if hits:
            return hits
    target = collapsed or needle
    words = page.get_text("words") or []
    if not words:
        return []
    from collections import defaultdict

    lines: dict = defaultdict(list)
    for item in words:
        x0, y0, x1, y1, word, block, line = item[:7]
        lines[(block, line)].append(item)
    rects = []
    for group in lines.values():
        group = sorted(group, key=lambda w: (w[0], w[1]))
        mapping = []
        compact_parts = []
        for w in group:
            piece = re.sub(r"\s+", "", w[4] or "")
            compact_parts.append(piece)
            mapping.extend([w] * len(piece))
        compact = "".join(compact_parts)
        start = 0
        while True:
            idx = compact.find(target, start)
            if idx < 0:
                break
            covered = mapping[idx: idx + len(target)]
            if covered:
                rects.append(fitz.Rect(
                    min(w[0] for w in covered),
                    min(w[1] for w in covered),
                    max(w[2] for w in covered),
                    max(w[3] for w in covered),
                ))
            start = idx + max(len(target), 1)
    return rects


def mask_pdf(input_path: Path, output_path: Path, replace_map: Dict[str, str]) -> None:
    """用红action 擦除原文并在原位置写入替换文本，输出仍为 PDF。"""
    import pymupdf as fitz

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(input_path))
    try:
        if getattr(doc, "is_encrypted", False):
            if not doc.authenticate(""):
                raise RuntimeError("加密 PDF 无法写入脱敏结果，请先解除密码保护")

        items = sorted(
            ((k, v) for k, v in replace_map.items() if k),
            key=lambda kv: len(kv[0]),
            reverse=True,
        )
        for page in doc:
            for widget in list(page.widgets() or []):
                value = widget.field_value
                if isinstance(value, str) and value:
                    new_val = apply_to_text(value, replace_map)
                    if new_val != value:
                        widget.field_value = new_val
                        widget.update()

            jobs = []
            for original, token in items:
                for rect in _pdf_search_rects(page, original):
                    jobs.append((rect, token))
            for rect, token in jobs:
                fontname = _pdf_fontname(token)
                fontsize = _pdf_fontsize_fit(rect, token, fontname)
                page.add_redact_annot(
                    rect,
                    text=token,
                    fontsize=fontsize,
                    fontname=fontname,
                    fill=(1, 1, 1),
                    text_color=(0, 0, 0),
                    cross_out=False,
                )
            if jobs:
                kwargs = {}
                if hasattr(fitz, "PDF_REDACT_IMAGE_NONE"):
                    kwargs["images"] = fitz.PDF_REDACT_IMAGE_NONE
                page.apply_redactions(**kwargs)

        doc.save(str(output_path), garbage=4, deflate=True)
    finally:
        doc.close()
